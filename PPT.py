#from pptx import Presentation  # try changing it to
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

# Create a presentation object
prs = Presentation()

# Slide 1: Introduction to SRE Automation
slide_1 = prs.slides.add_slide(prs.slide_layouts[5])

# Title
title_1 = slide_1.shapes.title
title_1.text = "Introduction to SRE Automation Solutions"

# Content
content_box_1 = slide_1.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(8.5), Inches(4))
content_frame_1 = content_box_1.text_frame
content_frame_1.text = "Brief overview of the purpose of SRE Automation: Enhancing reliability and efficiency in operations."
content_frame_1.add_paragraph().text = "Self-healing mechanisms for error management."
content_frame_1.add_paragraph().text = "Automation in device information retrieval during incidents."

# Visuals and Icons
content_frame_1.add_paragraph().text = "Icons: ELK, Jenkins, OpenShift, Icinga"

# Slide 2: Self-Healing Automation
slide_2 = prs.slides.add_slide(prs.slide_layouts[5])

# Title
title_2 = slide_2.shapes.title
title_2.text = "Self-Healing Automation: Responding to Errors Proactively"

# Content
content_box_2 = slide_2.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(8.5), Inches(4))
content_frame_2 = content_box_2.text_frame
content_frame_2.text = "Describe the trigger event: Continuous 503 or 500 errors detected."
content_frame_2.add_paragraph().text = "ELK Stack: Monitors and triggers a notification upon error detection."
content_frame_2.add_paragraph().text = "Jenkins Pipeline: Receives notification and communicates with OpenShift."
content_frame_2.add_paragraph().text = "OpenShift: Deletes the problematic pod automatically."
content_frame_2.add_paragraph().text = "Reduced manual intervention."
content_frame_2.add_paragraph().text = "Minimized performance disruptions and outages."

# Slide 3: Automation in Incident Management
slide_3 = prs.slides.add_slide(prs.slide_layouts[5])

# Title
title_3 = slide_3.shapes.title
title_3.text = "Enhancing Incident Response with Automated Device Info Retrieval"

# Content
content_box_3 = slide_3.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(8.5), Inches(4))
content_frame_3 = content_box_3.text_frame
content_frame_3.text = "Explain the role of Icinga in monitoring device statuses."
content_frame_3.add_paragraph().text = "Automatic fetching of device information by Icinga."
content_frame_3.add_paragraph().text = "Information availability boosts efficiency for SRE teams, ASCOE, and incident managers."
content_frame_3.add_paragraph().text = "Faster response times."
content_frame_3.add_paragraph().text = "Improved decision-making during crises."

# Slide 4: Summary & Key Takeaways
slide_4 = prs.slides.add_slide(prs.slide_layouts[5])

# Title
title_4 = slide_4.shapes.title
title_4.text = "Key Benefits and Takeaways"

# Content
content_box_4 = slide_4.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(8.5), Inches(4))
content_frame_4 = content_box_4.text_frame
content_frame_4.text = "Summarize the key points of both automation strategies."
content_frame_4.add_paragraph().text = "Significant reduction in manual tasks."
content_frame_4.add_paragraph().text = "Enhanced operational efficiency and reliability."
content_frame_4.add_paragraph().text = "Faster and more effective incident response."
content_frame_4.add_paragraph().text = "Encourage questions and discussions on potential implementation or integration."

# Save the presentation
prs.save("/Users/vinoji2005/Project/SRE-Automation/SRE_Automation_Presentation.pptx")
